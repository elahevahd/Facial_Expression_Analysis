from pptx import Presentation
from pptx.util import Inches
from pptx.util import Inches, Pt
import os 
import mean_plot_data_early_stop
import config
labeling= config.labeling

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Regression Network - Normalized Output"
subtitle.text = "Output normalized to [0,1]"

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
txBox = slide.shapes.add_textbox(Inches(3.5), Inches(0.5), Inches(3), Inches(3))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Configuration"
p.font.bold = True


for outer_fold in ['1']:
    for index in range(6):

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        img_path = 'Fold_{}/{}/plots/train_plot/{}.png'.format(outer_fold,labeling,index)
        pic = slide.shapes.add_picture(img_path,Inches(1.5) , Inches(3),width=Inches(7), height=Inches(3))

        txBox = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Training Results"
        p.font.size = Pt(15)

        txBox = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Fold {} - Parameter {}".format(outer_fold,index)
        p.font.bold = True


for outer_fold in ['2','3','4','5']:

    for index in range(2):

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        img_path = 'Fold_{}/{}/plots/train_plot/{}.png'.format(outer_fold,labeling,index)
        pic = slide.shapes.add_picture(img_path,Inches(1.5) , Inches(3),width=Inches(7), height=Inches(3))

        txBox = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Training Results"
        p.font.size = Pt(15)

        txBox = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Fold {} - Parameter {}".format(outer_fold,index)
        p.font.bold = True

# os.remove('train_results.pptx')
prs.save('train_results.pptx')
