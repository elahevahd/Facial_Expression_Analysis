from pptx import Presentation
from pptx.util import Inches
from pptx.util import Inches, Pt
import mean_plot_data_early_stop
import config
labeling= config.labeling

# params = {
# 'weight_decay': [1e-05] ,
# 'lr': [3e-03,5e-03,1e-02],
# 'dec_freq':[1],
# 'dropout_value':[0.3,0.5]
# }

params = {
'weight_decay': [1e-05] ,
'lr': [3e-03,5e-03],
'dec_freq':[1],
'dropout_value':[0.5]
}


number_of_pars=len(params['lr'])*len(params['weight_decay'])*len(params['dec_freq']*len(params['dropout_value']))



def slides_experiment(index_number):

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Regression Network VAS Results"
    subtitle.text = "Output normalized to [0,1]"

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    txBox = slide.shapes.add_textbox(Inches(3.5), Inches(0.5), Inches(3), Inches(3))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Configuration"
    p.font.bold = True

    for outer_fold in ['2','3','4','5']:

        for index in range(number_of_pars):

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            img_path = 'Fold_{}/{}/plots/loss/{}_outer.png'.format(outer_fold,labeling,index)
            pic = slide.shapes.add_picture(img_path,Inches(0.2),Inches(0.2), width=Inches(3), height=Inches(2))

            img_path = 'Fold_{}/{}/plots/scores/{}_outer_earlystop.png'.format(outer_fold,labeling,index)
            pic = slide.shapes.add_picture(img_path,Inches(3.5) , Inches(2.5),width=Inches(6), height=Inches(2))

            img_path = 'Fold_{}/{}/plots/scores/{}_dist_earlystop.png'.format(outer_fold,labeling,index)
            pic = slide.shapes.add_picture(img_path, Inches(3.5) , Inches(5) ,width=Inches(6), height=Inches(2))

            img_path = 'Fold_{}/{}/plots/distribution/{}_outer.png'.format(outer_fold,labeling,outer_fold)
            pic = slide.shapes.add_picture(img_path,Inches(0.5),Inches(2.5),width=Inches(3), height=Inches(2))

            img_path = 'Fold_{}/{}/plots/distribution/{}_dist.png'.format(outer_fold,labeling,outer_fold)
            pic = slide.shapes.add_picture(img_path,Inches(0.5),Inches(5),width=Inches(3), height=Inches(2))

            # (left, top, width, height)
            txBox = slide.shapes.add_textbox(Inches(6), Inches(1.75), Inches(2), Inches(1))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Extension Data Results"
            p.font.size = Pt(15)
            txBox = slide.shapes.add_textbox(Inches(6), Inches(4.25), Inches(2), Inches(1))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Distribution Data Results"
            p.font.size = Pt(15)
            txBox = slide.shapes.add_textbox(Inches(3.5), Inches(0.5), Inches(3), Inches(3))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Fold {} - Parameter {}".format(outer_fold,index)
            p.font.bold = True
            ######################### confusion matrix ######################
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            img_path = 'Fold_{}/{}/plots/confusion_matrix/confusion_{}.png'.format(outer_fold,labeling,index)
            pic = slide.shapes.add_picture(img_path,Inches(5) , Inches(1),width=Inches(4), height=Inches(3))
            img_path = 'Fold_{}/{}/plots/confusion_matrix/normed_confusion_{}.png'.format(outer_fold,labeling,index)
            pic = slide.shapes.add_picture(img_path,Inches(0.5),Inches(1),width=Inches(4), height=Inches(3))
            img_path = 'Fold_{}/{}/plots/confusion_matrix/dist_confusion_{}.png'.format(outer_fold,labeling,index)
            pic = slide.shapes.add_picture(img_path, Inches(5) , Inches(4) ,width=Inches(4), height=Inches(3))
            img_path = 'Fold_{}/{}/plots/confusion_matrix/dist_normed_confusion_{}.png'.format(outer_fold,labeling,index)
            pic = slide.shapes.add_picture(img_path,Inches(0.5),Inches(4),width=Inches(4), height=Inches(3))

            # (left, top, width, height)
            txBox = slide.shapes.add_textbox(Inches(4), Inches(0.5), Inches(2), Inches(1))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Extension Data Results"
            p.font.size = Pt(12)
            txBox = slide.shapes.add_textbox(Inches(4), Inches(3.5), Inches(2), Inches(1))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Distribution Data Results"
            p.font.size = Pt(12)
            txBox = slide.shapes.add_textbox(Inches(3.5), Inches(0.05), Inches(3), Inches(3))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Fold {} - Parameter {}".format(outer_fold,index)
            p.font.bold = True

    # blank_slide_layout = prs.slide_layouts[6]
    # slide = prs.slides.add_slide(blank_slide_layout)
    # img_path = 'mean_results/{}/mean_extension_fig.png'.format(labeling)
    # pic = slide.shapes.add_picture(img_path,Inches(0.2) , Inches(2),width=Inches(5), height=Inches(3))
    # img_path = 'VAS_Extension.png'
    # pic = slide.shapes.add_picture(img_path,Inches(5.5) , Inches(2),width=Inches(4), height=Inches(3))

    # blank_slide_layout = prs.slide_layouts[6]
    # slide = prs.slides.add_slide(blank_slide_layout)
    # img_path = 'mean_results/{}/mean_distribution_fig.png'.format(labeling)
    # pic = slide.shapes.add_picture(img_path,Inches(0.2) , Inches(2),width=Inches(5), height=Inches(3))
    # img_path = 'VAS_Distribution.png'
    # pic = slide.shapes.add_picture(img_path,Inches(5.5) , Inches(2),width=Inches(4), height=Inches(3))

    prs.save(labeling+'_results_{}.pptx'.format(index_number))




slides_experiment('dropout')